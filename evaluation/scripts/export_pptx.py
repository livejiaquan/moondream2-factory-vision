#!/usr/bin/env python3
"""
Export HTML presentation to PPTX by screenshotting each slide.

Usage:
  conda run -n moondream python scripts/export_pptx.py --run 20260421T100247Z
"""
from __future__ import annotations

import argparse
import io
from pathlib import Path

EVAL_DIR    = Path(__file__).resolve().parents[1]
OUTPUTS_DIR = EVAL_DIR / "outputs"

SLIDE_W_PX = 1280
SLIDE_H_PX = 720
# PPTX slide dimensions in EMUs (English Metric Units): 1 inch = 914400 EMU
# 1280x720 at 96 DPI → 13.33 x 7.5 inches
PPTX_W = int(13.333 * 914400)
PPTX_H = int(7.5    * 914400)


def screenshot_slides(html_path: Path) -> list[bytes]:
    from playwright.sync_api import sync_playwright
    images: list[bytes] = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": SLIDE_W_PX, "height": SLIDE_H_PX})
        page.goto(html_path.as_uri(), wait_until="networkidle")
        slide_count = page.evaluate("document.querySelectorAll('.slide').length")
        print(f"[screenshot] {slide_count} slides found")
        for i in range(slide_count):
            page.evaluate(f"""
                document.querySelectorAll('.slide').forEach((s, j) => {{
                    s.style.display = j === {i} ? 'flex' : 'none';
                }});
            """)
            page.wait_for_timeout(150)
            png = page.screenshot(full_page=False)
            images.append(png)
            print(f"[screenshot] slide {i+1}/{slide_count}")
        browser.close()
    return images


def build_pptx(images: list[bytes], out_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width  = Emu(PPTX_W)
    prs.slide_height = Emu(PPTX_H)

    blank_layout = prs.slide_layouts[6]  # completely blank

    for i, png_bytes in enumerate(images):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(png_bytes),
            left=Emu(0), top=Emu(0),
            width=Emu(PPTX_W), height=Emu(PPTX_H),
        )
        print(f"[pptx] added slide {i+1}/{len(images)}")

    prs.save(str(out_path))
    print(f"[pptx] saved → {out_path}")


def main() -> int:
    p = argparse.ArgumentParser()
    p.add_argument("--run", required=True)
    args = p.parse_args()

    run_dir  = OUTPUTS_DIR / args.run
    html_path = run_dir / "presentation.html"
    out_path  = run_dir / "presentation.pptx"

    if not html_path.exists():
        print(f"[error] {html_path} not found — run generate_presentation.py first")
        return 1

    print(f"[export] {html_path}")
    images = screenshot_slides(html_path)
    build_pptx(images, out_path)
    print(f"\n完成：{out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
