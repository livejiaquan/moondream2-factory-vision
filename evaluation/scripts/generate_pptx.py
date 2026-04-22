#!/usr/bin/env python3
"""
Generate editable PPTX from evaluation results.

Usage:
  conda run -n moondream python scripts/generate_pptx.py --run 20260421T100247Z
"""
from __future__ import annotations

import argparse
import io
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[2]
EVAL_DIR     = Path(__file__).resolve().parents[1]
OUTPUTS_DIR  = EVAL_DIR / "outputs"

# ── Design canvas: 1280×720 ────────────────────────────────────────────────
# Scale: 1280px → 9144000 EMU (10 inch wide PowerPoint slide)
_SCALE = 9144000 / 1280

def px(n: float) -> int:
    return int(n * _SCALE)

def rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

# ── Colour palette ─────────────────────────────────────────────────────────
C = {
    "bg":       "#F8FAFC",
    "card":     "#FFFFFF",
    "card2":    "#F1F5F9",
    "teal":     "#0D9488",
    "blue":     "#2563EB",
    "amber":    "#D97706",
    "red":      "#DC2626",
    "text":     "#0F172A",
    "sub":      "#334155",
    "muted":    "#64748B",
    "border":   "#E2E8F0",
    "hdr_bg":   "#0F172A",
    "hdr_text": "#FFFFFF",
}


# ── Shape helpers ──────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill: str | None = None,
             line: str | None = None, line_pt: float = 1.5,
             radius: bool = False) -> object:
    st = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(st, px(x), px(y), px(w), px(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_img(slide, img_bytes: bytes, x, y, w, h) -> None:
    slide.shapes.add_picture(io.BytesIO(img_bytes), px(x), px(y), px(w), px(h))


def add_text(slide, x, y, w, h, text: str, size: float, color: str,
             bold: bool = False, align: PP_ALIGN = PP_ALIGN.LEFT,
             wrap: bool = True, italic: bool = False,
             margin_l: float = 4, margin_t: float = 2) -> object:
    txBox = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.margin_left  = px(margin_l)
    tf.margin_right = px(margin_l)
    tf.margin_top   = px(margin_t)
    tf.margin_bottom = px(margin_t)
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return txBox


def add_multiline(slide, x, y, w, h, lines: list[tuple[str, float, str, bool]],
                  align: PP_ALIGN = PP_ALIGN.LEFT,
                  line_spacing: float = 1.2) -> None:
    """lines = [(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = px(6)
    tf.margin_right  = px(4)
    tf.margin_top    = px(4)
    tf.margin_bottom = px(4)
    for i, (text, size, color, bold) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        para.alignment = align


def hdr(slide, title: str, subtitle: str = "") -> None:
    """Dark header bar across the top."""
    add_rect(slide, 0, 0, 1280, 56, fill=C["hdr_bg"])
    add_rect(slide, 0, 53, 1280, 3, fill=C["teal"])  # teal accent line
    add_text(slide, 14, 6, 900, 30, title, 20, C["hdr_text"], bold=True)
    if subtitle:
        add_text(slide, 14, 32, 1100, 20, subtitle, 10, C["teal"], bold=True)


def card_rect(slide, x, y, w, h, accent_color: str | None = None,
              fill: str = "card", accent_side: str = "top",
              radius: bool = True) -> None:
    add_rect(slide, x, y, w, h, fill=C[fill], radius=radius)
    if accent_color:
        if accent_side == "top":
            add_rect(slide, x, y, w, 5, fill=accent_color)
        elif accent_side == "left":
            add_rect(slide, x, y, 5, h, fill=accent_color)


# ── Frame extraction ───────────────────────────────────────────────────────
_video_cache: dict[str, Path | None] = {}
_frame_cache: dict[tuple, bytes | None] = {}

def _resolve_video(name: str, data_dir: Path) -> Path | None:
    if name in _video_cache:
        return _video_cache[name]
    p = data_dir / name
    _video_cache[name] = p if p.exists() else next(data_dir.rglob(name), None)
    return _video_cache[name]


def get_frame_bytes(video: str, time_sec: float, data_dir: Path,
                    width: int = 480) -> bytes | None:
    key = (video, time_sec, width)
    if key in _frame_cache:
        return _frame_cache[key]
    try:
        import cv2
        from PIL import Image
        vp = _resolve_video(video, data_dir)
        if vp is None:
            _frame_cache[key] = None
            return None
        cap = cv2.VideoCapture(str(vp))
        cap.set(cv2.CAP_PROP_POS_MSEC, float(time_sec) * 1000.0)
        ok, frame = cap.read()
        cap.release()
        if not ok or frame is None:
            _frame_cache[key] = None
            return None
        img = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
        ratio = width / img.width
        img = img.resize((width, int(img.height * ratio)), Image.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, "JPEG", quality=82)
        result = buf.getvalue()
        _frame_cache[key] = result
        return result
    except Exception:
        _frame_cache[key] = None
        return None


def placeholder_img(w_px: int, h_px: int) -> bytes:
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (w_px, h_px), color=(241, 245, 249))
    d = ImageDraw.Draw(img)
    d.text((w_px // 2 - 30, h_px // 2 - 10), "影片不存在", fill=(100, 116, 139))
    buf = io.BytesIO()
    img.save(buf, "JPEG", quality=70)
    return buf.getvalue()


# ── Data loaders ───────────────────────────────────────────────────────────
def _load_phase2(run: str) -> list[dict]:
    p = OUTPUTS_DIR / f"{run}_occlusion_variants" / "phase2_full.jsonl"
    return [json.loads(l) for l in p.read_text("utf-8").splitlines() if l.strip()]


def _load_mask_data(run: str):
    d = OUTPUTS_DIR / f"{run}_occlusion_variants"
    phase1 = {(r["video"], float(r["time_sec"])): r
              for l in (d / "mask_phase2_full.jsonl").read_text().splitlines()
              if l.strip() for r in [json.loads(l)]}
    aof = {(r["video"], float(r["time_sec"])): r
           for l in (d / "mask_fullrun_anything_on_face.jsonl").read_text().splitlines()
           if l.strip() for r in [json.loads(l)]}
    fp_keys = [k for k, r in phase1.items() if r["baseline"]["missing_face_mask"] == "yes"]
    ok_keys = [k for k, r in phase1.items() if r["baseline"]["missing_face_mask"] == "no"]
    return phase1, aof, fp_keys, ok_keys


def _load_baseline(run: str) -> list[dict]:
    p = OUTPUTS_DIR / run / "results.jsonl"
    return [json.loads(l) for l in p.read_text("utf-8").splitlines() if l.strip()]


# ── Slides ─────────────────────────────────────────────────────────────────
BODY_Y  = 62   # px below header
BODY_H  = 720 - BODY_Y - 12
BODY_X  = 16
BODY_W  = 1280 - BODY_X * 2


def s_cover(prs: Presentation, run: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1280, 720, fill=C["bg"])

    add_rect(slide, 0, 0, 400, 720, fill=C["hdr_bg"])
    add_rect(slide, 0, 0, 6, 720, fill=C["teal"])

    add_rect(slide, 22, 24, 356, 28, fill=C["teal"], radius=True)
    add_text(slide, 26, 26, 348, 24, "USIG · MOONDREAM EVALUATION",
             9, "#FFFFFF", bold=True, align=PP_ALIGN.CENTER)

    add_multiline(slide, 22, 62, 356, 180, [
        ("Moondream", 32, "#FFFFFF", True),
        ("PPE 偵測驗證報告", 28, "#FFFFFF", True),
    ])

    add_multiline(slide, 22, 250, 356, 160, [
        ("Helmet & Mask 辨識效果", 14, "#94a3b8", False),
        ("Occlusion 實驗", 14, "#94a3b8", False),
        ("Prompt 策略比較", 14, "#94a3b8", False),
    ])
    add_text(slide, 22, 420, 356, 24, f"{run} · Moondream2 2025-01-09",
             11, "#64748b")

    add_text(slide, 430, 80, 820, 24, "本次實驗規模",
             11, C["muted"], bold=True, align=PP_ALIGN.CENTER)

    stats = [("7", "CCTV 影片"), ("47", "抽樣 Frames"),
             ("2", "PPE 目標"), ("3", "實驗階段")]
    cols_x = [430, 720]
    rows_y = [114, 380]
    for idx, (val, lbl) in enumerate(stats):
        cx = cols_x[idx % 2]
        cy = rows_y[idx // 2]
        card_rect(slide, cx, cy, 250, 230, accent_color=C["teal"])
        add_text(slide, cx + 10, cy + 30, 230, 100, val,
                 52, C["teal"], bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, cx + 10, cy + 145, 230, 60, lbl,
                 15, C["sub"], align=PP_ALIGN.CENTER)


def s_fp_overview(prs: Presentation, run: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1280, 720, fill=C["bg"])
    hdr(slide, "有人幀 FP 率總覽", "FALSE POSITIVE RATE · PERSON-GATED FRAMES (n=34)")

    rows = [
        ("安全帽缺失 (missing_helmet)", 32.4, "11/34", C["amber"]),
        ("口罩缺失 (missing_face_mask)", 70.6, "24/34", C["red"]),
        ("抽菸偵測 (smoking_visible)",   8.8,  " 3/34", C["teal"]),
        ("異常行為 (abnormal_behavior)", 8.8,  " 3/34", C["teal"]),
        ("倒地偵測 (fallen_person)",     0.0,  " 0/34", C["teal"]),
    ]

    bar_x, bar_w = 290, 820
    for i, (name, pct, frac, color) in enumerate(rows):
        y = BODY_Y + 18 + i * 80
        add_text(slide, BODY_X + 4, y + 4, 270, 28, name, 13, C["text"])
        add_rect(slide, bar_x, y, bar_w, 34, fill=C["card2"])
        fill_w = max(4, int(bar_w * pct / 100))
        add_rect(slide, bar_x, y, fill_w, 34, fill=color)
        add_text(slide, bar_x + 8, y + 8, 200, 20, f"{frac} 幀",
                 12, "#FFFFFF", bold=True)
        add_text(slide, bar_x + bar_w + 8, y + 6, 70, 24,
                 f"{pct:.1f}%", 14, color, bold=True)

    # summary cards
    summaries = [
        ("最高優先改善", "Mask 70.6% FP", "幾乎等於隨機，不可直接用於告警", C["red"]),
        ("中等優先改善", "Helmet 32.4% FP", "遮擋/角度問題導致，已有解法", C["amber"]),
        ("已達可用標準", "其他 Query ≤ 8.8%", "加交叉驗證規則後可直接整合", C["teal"]),
    ]
    for j, (sub, title, desc, color) in enumerate(summaries):
        cx = BODY_X + j * 415
        cy = BODY_Y + 440
        card_rect(slide, cx, cy, 400, 200, accent_color=color, accent_side="left")
        add_text(slide, cx + 14, cy + 8, 375, 22, sub, 11, C["muted"])
        add_text(slide, cx + 14, cy + 32, 375, 28, title, 15, C["text"], bold=True)
        add_text(slide, cx + 14, cy + 66, 375, 50, desc, 13, C["sub"], wrap=True)


def s_helmet_fp_causes(prs: Presentation, phase2_rows: list[dict], data_dir: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1280, 720, fill=C["bg"])
    hdr(slide, "Helmet FP 原因分析 — 遮擋 / 角度問題",
        "HELMET · DETECT=0 FRAMES · helmet_area_visible EXPERIMENT (n=7)")

    occ  = [r for r in phase2_rows if r["baseline"]["detect_helmet"] == 0]
    good = [r for r in occ if r["can_confirm_helmet"]["label"] == "no"][:2]
    miss = [r for r in occ if r["can_confirm_helmet"]["label"] == "yes"][:1]

    # left stat cards
    for i, (num, lbl, color) in enumerate([
        ("6/7", "detect=0 幀\n正確識別為遮擋", C["teal"]),
        ("1/7", "識別失誤\n（不應取消告警）", C["red"]),
    ]):
        cy = BODY_Y + 8 + i * 150
        card_rect(slide, BODY_X, cy, 220, 134, accent_color=color, accent_side="left")
        add_text(slide, BODY_X + 14, cy + 12, 200, 64, num,
                 40, color, bold=True)
        add_text(slide, BODY_X + 14, cy + 78, 200, 46, lbl,
                 12, C["sub"], wrap=True)

    card_rect(slide, BODY_X, BODY_Y + 316, 220, 340, fill="card2")
    add_text(slide, BODY_X + 10, BODY_Y + 324, 200, 20,
             "PIPELINE", 9, C["muted"], bold=True)
    add_multiline(slide, BODY_X + 10, BODY_Y + 350, 200, 290, [
        ("detect=0", 13, C["sub"], False),
        ("↓ helmet_area_visible", 13, C["sub"], False),
        ("= no  → 取消告警", 13, C["red"], True),
        ("= yes → 升級告警 ⚠", 13, C["teal"], True),
    ])

    # right 2×2 grid
    imgs = good[:2] + (miss if miss else [good[0] if good else None])
    # fill to 4
    while len(imgs) < 4 and good:
        imgs.append(good[0])
    imgs = imgs[:4]

    grid_x, grid_y = 252, BODY_Y + 8
    iw, ih = 490, 318
    gap = 12
    labels_map = {
        "no":  ("遮擋識別成功 ✓", C["teal"]),
        "yes": ("識別失誤 ✗",     C["red"]),
    }
    for idx, r in enumerate(imgs):
        if r is None:
            continue
        gx = grid_x + (idx % 2) * (iw + gap)
        gy = grid_y + (idx // 2) * (ih + gap)
        label = r["can_confirm_helmet"]["label"]
        lbl_text, lbl_color = labels_map.get(label, ("—", C["muted"]))
        fb = get_frame_bytes(r["video"], float(r["time_sec"]), data_dir, width=600)
        if fb is None:
            fb = placeholder_img(600, 338)
        add_rect(slide, gx, gy, iw, ih, fill=C["card2"])
        add_img(slide, fb, gx, gy, iw, ih - 52)
        add_rect(slide, gx, gy + ih - 52, iw, 52, fill=C["hdr_bg"])
        add_text(slide, gx + 6, gy + ih - 48, iw - 12, 22,
                 f"helmet_area_visible = {label}   {lbl_text}",
                 11, lbl_color, bold=True)
        add_text(slide, gx + 6, gy + ih - 26, iw - 12, 20,
                 f"{Path(r['video']).name[:40]} @ {r['time_sec']}s",
                 10, "#cbd5e1")


def s_helmet_contradiction(prs: Presentation, phase2_rows: list[dict], data_dir: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1280, 720, fill=C["bg"])
    hdr(slide, "Detect vs Query 矛盾案例",
        "HELMET · DETECT=1 BUT QUERY SAYS MISSING · n=4 FRAMES")

    contra = [r for r in phase2_rows
              if r["baseline"]["detect_helmet"] > 0
              and r["baseline"]["missing_helmet"] == "yes"][:4]

    # info cards top
    card_rect(slide, BODY_X, BODY_Y + 4, 820, 120, accent_color=C["blue"], accent_side="left")
    add_text(slide, BODY_X + 16, BODY_Y + 14, 800, 28,
             "現象：YOLO detect 到安全帽，但 Moondream query 說工人缺少安全帽", 14, C["text"], bold=True)
    add_text(slide, BODY_X + 16, BODY_Y + 48, 800, 60,
             "helmet_area_visible 全部回答 yes（能看清楚頭部）\n→ 確認 detect 比 query 更可靠，query 有誤報傾向",
             13, C["sub"], wrap=True)

    card_rect(slide, BODY_X + 844, BODY_Y + 4, 404, 120, accent_color=C["teal"], accent_side="left")
    add_multiline(slide, BODY_X + 860, BODY_Y + 14, 380, 100, [
        ("實驗結論", 11, C["muted"], False),
        ("detect 出現時 → 信任 detect", 13, C["sub"], False),
        ("detect = 0 → 才問 query", 13, C["sub"], False),
        ("query 優先級低於 detect", 13, C["sub"], False),
    ])

    add_text(slide, BODY_X, BODY_Y + 134, 600, 18,
             "矛盾幀抽樣（detect 框到帽但 query 誤報缺少）", 10, C["muted"], bold=True)

    # 2×2 image grid
    iw, ih, gap = 608, 240, 12
    for idx, r in enumerate(contra):
        gx = BODY_X + (idx % 2) * (iw + gap)
        gy = BODY_Y + 156 + (idx // 2) * (ih + gap)
        label = r["can_confirm_helmet"]["label"]
        fb = get_frame_bytes(r["video"], float(r["time_sec"]), data_dir, width=640)
        if fb is None:
            fb = placeholder_img(640, 360)
        add_rect(slide, gx, gy, iw, ih, fill=C["card2"])
        add_img(slide, fb, gx, gy, iw, ih - 52)
        add_rect(slide, gx, gy + ih - 52, iw, 52, fill=C["hdr_bg"])
        add_text(slide, gx + 6, gy + ih - 49, iw - 12, 22,
                 f"detect=1  query=missing  helmet_area_visible={label}",
                 11, C["teal"], bold=True)
        add_text(slide, gx + 6, gy + ih - 27, iw - 12, 20,
                 f"{Path(r['video']).name[:45]} @ {r['time_sec']}s",
                 10, "#cbd5e1")


def s_mask_problem(prs: Presentation, mask_data: tuple) -> None:
    phase1, aof, fp_keys, ok_keys = mask_data
    fav_fp_no  = sum(1 for k in fp_keys if phase1[k]["face_area_visible"]["label"] == "no")
    fav_ok_no  = sum(1 for k in ok_keys if phase1[k]["face_area_visible"]["label"] == "no")
    aof_fp_yes = sum(1 for k in fp_keys if aof.get(k, {}).get("anything_on_face", {}).get("label") == "yes")
    aof_ok_yes = sum(1 for k in ok_keys if aof.get(k, {}).get("anything_on_face", {}).get("label") == "yes")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1280, 720, fill=C["bg"])
    hdr(slide, "Mask 偵測：兩種方法對比",
        "MASK · face_area_visible vs anything_on_face")

    for col, (accent, method_lbl, q1_num, q1_lbl, q2_num, q2_lbl,
              sign, insight) in enumerate([
        (C["red"], "方法一（失敗）",
         f"{fav_fp_no}/{len(fp_keys)}", "YOLO 告警幀\nface_area_visible=no",
         f"{fav_ok_no}/{len(ok_keys)}", "YOLO 正常幀\nface_area_visible=no",
         "=", "兩組結果相同 → 無法區分誤報與正常幀"),
        (C["teal"], "方法二（有效）",
         f"{aof_fp_yes}/{len(fp_keys)}", "YOLO 告警幀\nanything_on_face=yes",
         f"{aof_ok_yes}/{len(ok_keys)}", "YOLO 正常幀\nanything_on_face=yes",
         "≠", "兩組結果不同 → FP 告警率從 70.6% 降到 26.5%"),
    ]):
        cx = BODY_X + col * 500
        cw = 488
        card_rect(slide, cx, BODY_Y + 6, cw, 622, accent_color=accent)
        add_text(slide, cx + 12, BODY_Y + 20, cw - 20, 26,
                 method_lbl, 15, accent, bold=True)
        add_text(slide, cx + 12, BODY_Y + 54, cw - 20, 30,
                 "問「能不能看到臉部？」" if col == 0 else "問「臉上有沒有防護裝備？」",
                 14, C["text"], bold=True)

        # two stat boxes + sign
        for si, (num, lbl, ncolor) in enumerate([
            (q1_num, q1_lbl, accent),
            (q2_num, q2_lbl, C["blue"] if col == 1 else accent),
        ]):
            sx = cx + 12 + si * 210
            sy = BODY_Y + 96
            add_rect(slide, sx, sy, 195, 200, fill=C["card2"], radius=True)
            add_text(slide, sx + 8, sy + 18, 179, 80,
                     num, 36, ncolor, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, sx + 8, sy + 108, 179, 74,
                     lbl, 12, C["sub"], align=PP_ALIGN.CENTER, wrap=True)

        add_text(slide, cx + 215, BODY_Y + 166, 50, 40,
                 sign, 26, C["muted"], bold=True, align=PP_ALIGN.CENTER)

        # insight box
        add_rect(slide, cx + 12, BODY_Y + 308, cw - 24, 70,
                 fill=C["card2"], radius=True)
        add_rect(slide, cx + 12, BODY_Y + 308, 5, 70, fill=accent)
        add_text(slide, cx + 22, BODY_Y + 322, cw - 40, 48,
                 insight, 13, accent, bold=True, wrap=True)

    # Root cause box (right column)
    rx = BODY_X + 1010
    card_rect(slide, rx, BODY_Y + 6, 238, 622, accent_color=C["amber"])
    add_text(slide, rx + 12, BODY_Y + 20, 214, 26, "根本原因", 14, C["amber"], bold=True)
    add_text(slide, rx + 12, BODY_Y + 52, 214, 56,
             "俯角攝影機\n+ 帽沿遮擋下臉", 14, C["text"], bold=True, wrap=True)

    for by, (lbl, desc, bcolor) in enumerate([
        ("Helmet", "頭頂最清楚，問「能不能看到」有效", C["teal"]),
        ("Mask",   "下臉被帽沿遮住，改問「有沒有防護裝備」", C["amber"]),
        ("",       "不是模型看不到人，而是原本問題設計錯了。", C["red"]),
    ]):
        by2 = BODY_Y + 126 + by * 150
        add_rect(slide, rx + 12, by2, 214, 134, fill=C["card2"], radius=True)
        add_rect(slide, rx + 12, by2, 5, 134, fill=bcolor)
        if lbl:
            add_text(slide, rx + 22, by2 + 8, 190, 24, lbl, 13, bcolor, bold=True)
        add_text(slide, rx + 22, by2 + (32 if lbl else 12), 190, 90,
                 desc, 12, C["sub"], wrap=True)


def s_mask_fp_photos(prs: Presentation, mask_data: tuple, data_dir: Path) -> None:
    phase1, aof, fp_keys, _ = mask_data
    sample = fp_keys[:4]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1280, 720, fill=C["bg"])
    hdr(slide, "Mask 誤報幀實際照片（anything_on_face 驗證結果）",
        "FP FRAMES · GREEN=FACE COVERING CONFIRMED→ALERT CANCELLED · RED=UNCONFIRMED→ALERT KEPT")

    iw, ih, gap = 612, 318, 12
    for idx, k in enumerate(sample):
        r = phase1[k]
        fav  = phase1[k]["face_area_visible"]["label"]
        aof_lbl = aof.get(k, {}).get("anything_on_face", {}).get("label", "—")
        detected = aof_lbl == "yes"
        bcolor = C["teal"] if detected else C["red"]
        gx = BODY_X + (idx % 2) * (iw + gap)
        gy = BODY_Y + 8 + (idx // 2) * (ih + gap)
        fb = get_frame_bytes(r["video"], float(r["time_sec"]), data_dir, width=640)
        if fb is None:
            fb = placeholder_img(640, 360)
        add_rect(slide, gx, gy, iw, ih, fill=C["card2"])
        add_img(slide, fb, gx, gy, iw, ih - 58)
        add_rect(slide, gx, gy + ih - 58, iw, 58, fill=C["hdr_bg"])
        add_rect(slide, gx, gy, iw, 4, fill=bcolor)
        result_txt = "anything_on_face=yes → 告警取消" if detected else "anything_on_face=no → 告警保留"
        add_text(slide, gx + 6, gy + ih - 54, iw - 12, 22,
                 f"face_area_visible={fav}   {result_txt}",
                 10, bcolor, bold=True)
        add_text(slide, gx + 6, gy + ih - 30, iw - 12, 24,
                 f"{Path(r['video']).name[:42]} @ {r['time_sec']}s",
                 10, "#cbd5e1")


def s_mask_improvement_stats(prs: Presentation, mask_data: tuple) -> None:
    phase1, aof, fp_keys, ok_keys = mask_data
    aof_fp_yes = sum(1 for k in fp_keys if aof.get(k, {}).get("anything_on_face", {}).get("label") == "yes")
    aof_ok_yes = sum(1 for k in ok_keys if aof.get(k, {}).get("anything_on_face", {}).get("label") == "yes")
    aof_fp_no  = len(fp_keys) - aof_fp_yes

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1280, 720, fill=C["bg"])
    hdr(slide, "Mask — anything_on_face 全量結果",
        "FULL RUN · n=34 PERSON FRAMES · ALERT RATE 70.6% → 26.5%")

    # 3 big stat cards
    stats = [
        (str(aof_fp_yes) + "/" + str(len(fp_keys)),
         "YOLO: missing_face_mask=yes\nMoondream: anything_on_face=yes\n→ 告警取消",
         C["teal"]),
        (str(aof_ok_yes) + "/" + str(len(ok_keys)),
         "YOLO: missing_face_mask=no\nMoondream: anything_on_face=yes\n→ 兩者一致，無告警",
         C["blue"]),
        (str(aof_fp_no) + "/" + str(len(fp_keys)),
         "YOLO: missing_face_mask=yes\nMoondream: anything_on_face=no\n→ 告警保留（仍為 FP）",
         C["red"]),
    ]
    sw = 296
    for i, (num, desc, color) in enumerate(stats):
        cx = BODY_X + i * (sw + 12)
        card_rect(slide, cx, BODY_Y + 8, sw, 310, accent_color=color)
        add_text(slide, cx + 12, BODY_Y + 24, sw - 24, 100,
                 num, 52, color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, cx + 12, BODY_Y + 136, sw - 24, 170,
                 desc, 12, C["sub"], wrap=True)

    # Bar comparison
    bar_y = BODY_Y + 334
    card_rect(slide, BODY_X, bar_y, 930, 200)
    add_text(slide, BODY_X + 16, bar_y + 12, 500, 20,
             "告警率比較（全部人員幀 n=34）", 11, C["muted"], bold=True)
    bar_x, bar_w = BODY_X + 16, 700
    for bi, (label, pct, color) in enumerate([
        ("face_area_visible (無效)", 70.6, C["red"]),
        ("anything_on_face (有效)", 26.5, C["teal"]),
    ]):
        by = bar_y + 44 + bi * 72
        add_text(slide, bar_x, by - 2, 380, 20, label, 12, C["sub"])
        add_rect(slide, bar_x, by + 20, bar_w, 30, fill=C["card2"])
        fw = max(4, int(bar_w * pct / 100))
        add_rect(slide, bar_x, by + 20, fw, 30, fill=color)
        add_text(slide, bar_x + bar_w + 8, by + 22, 80, 26,
                 f"{pct}%", 16, color, bold=True)

    # Limitation box
    card_rect(slide, BODY_X + 958, BODY_Y + 8, 290, 526,
              accent_color=C["amber"])
    add_text(slide, BODY_X + 972, BODY_Y + 20, 264, 22,
             "殘留 FP 原因", 11, C["amber"], bold=True)
    add_text(slide, BODY_X + 972, BODY_Y + 48, 264, 30,
             "9/24 幀仍無法確認", 16, C["text"], bold=True)
    add_multiline(slide, BODY_X + 972, BODY_Y + 90, 264, 280, [
        ("→ 工人背對鏡頭", 13, C["sub"], False),
        ("→ 臉部完全不可見", 13, C["sub"], False),
        ("→ 非 Prompt 問題", 13, C["sub"], False),
        ("→ 鏡頭角度死角", 13, C["sub"], False),
    ])
    add_rect(slide, BODY_X + 972, BODY_Y + 400, 264, 80, fill="#fef3c7", radius=True)
    add_text(slide, BODY_X + 980, BODY_Y + 412, 248, 60,
             "硬體層面問題，需補充側面鏡頭覆蓋", 12, "#92400e", wrap=True)


def _image_grid_slide(prs, title, subtitle, keys, phase1, aof, data_dir,
                      tag_fn, bcolor_fn, cols=3) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1280, 720, fill=C["bg"])
    hdr(slide, title, subtitle)

    n = len(keys)
    rows = 2
    gap = 10
    iw = (BODY_W - gap * (cols - 1)) // cols
    ih = (BODY_H - gap * (rows - 1)) // rows

    for idx, k in enumerate(keys[:cols * rows]):
        r = phase1[k]
        gx = BODY_X + (idx % cols) * (iw + gap)
        gy = BODY_Y + 8 + (idx // cols) * (ih + gap)
        bc = bcolor_fn(k, aof)
        tag = tag_fn(k, aof)
        fb = get_frame_bytes(r["video"], float(r["time_sec"]), data_dir, width=440)
        if fb is None:
            fb = placeholder_img(440, 248)
        add_rect(slide, gx, gy, iw, ih, fill=C["card2"])
        add_img(slide, fb, gx, gy, iw, ih - 46)
        add_rect(slide, gx, gy + ih - 46, iw, 46, fill=C["hdr_bg"])
        add_rect(slide, gx, gy, iw, 4, fill=bc)
        add_text(slide, gx + 5, gy + ih - 43, iw - 10, 22,
                 tag, 10, bc, bold=True)
        add_text(slide, gx + 5, gy + ih - 22, iw - 10, 18,
                 f"{Path(r['video']).name[:30]} @ {r['time_sec']}s",
                 9, "#cbd5e1")


def s_mask_detected_frames(prs: Presentation, mask_data: tuple, data_dir: Path) -> None:
    phase1, aof, fp_keys, _ = mask_data
    keys = [k for k in fp_keys if aof.get(k, {}).get("anything_on_face", {}).get("label") == "yes"][:6]
    _image_grid_slide(
        prs,
        "Mask — YOLO 告警且 anything_on_face = yes 的幀（告警取消）",
        f"missing_face_mask=yes (YOLO) → anything_on_face=yes (Moondream) → alert cancelled · {len(keys)}/24",
        keys, phase1, aof, data_dir,
        tag_fn=lambda k, _aof: "anything_on_face = yes → 告警取消",
        bcolor_fn=lambda k, _aof: C["teal"],
        cols=3,
    )


def s_mask_missed_frames(prs: Presentation, mask_data: tuple, data_dir: Path) -> None:
    phase1, aof, fp_keys, _ = mask_data
    keys = [k for k in fp_keys if aof.get(k, {}).get("anything_on_face", {}).get("label") != "yes"][:6]
    _image_grid_slide(
        prs,
        "Mask — YOLO 告警且 anything_on_face = no 的幀（告警保留）",
        "missing_face_mask=yes (YOLO) → anything_on_face=no (Moondream) → alert kept · 9/24 FP FRAMES REMAIN",
        keys, phase1, aof, data_dir,
        tag_fn=lambda k, _aof: "anything_on_face = no → 告警保留",
        bcolor_fn=lambda k, _aof: C["red"],
        cols=3,
    )


def s_per_video_breakdown(prs: Presentation, baseline: list[dict]) -> None:
    rows = baseline
    person = [r for r in rows if r["result"]["detections"]["person"]["count"] > 0]
    vids = sorted(set(r["video"] for r in rows))
    queries = [
        ("missing_helmet",    "missing_helmet",    C["amber"]),
        ("missing_face_mask", "missing_face_mask", C["red"]),
        ("abnormal_behavior", "abnormal_behavior", C["blue"]),
    ]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1280, 720, fill=C["bg"])
    hdr(slide, "各影片誤報分布",
        "PER-VIDEO FALSE POSITIVE BREAKDOWN · PERSON-GATED FRAMES")

    # legend
    lx = BODY_X
    for i, (key, label, color) in enumerate(queries):
        add_rect(slide, lx + i * 310, BODY_Y + 10, 14, 14, fill=color)
        add_text(slide, lx + i * 310 + 18, BODY_Y + 8, 280, 18, label, 11, C["sub"])

    bar_y_start = BODY_Y + 38
    row_h = 138
    bar_max_w = 820
    name_w = 310

    for vi, vid in enumerate(vids):
        vrows = [r for r in person if r["video"] == vid]
        n = len(vrows)
        if n == 0:
            continue
        ry = bar_y_start + vi * row_h
        short = Path(vid).stem[:38]
        add_text(slide, BODY_X, ry + 4, name_w - 8, 22,
                 short, 11, C["text"], bold=True)
        add_text(slide, BODY_X, ry + 24, name_w - 8, 18,
                 f"n={n} 有人幀", 10, C["muted"])

        for bi, (key, _, color) in enumerate(queries):
            yes = sum(1 for r in vrows if r["result"]["queries"].get(key, {}).get("label") == "yes")
            pct = yes / n
            bx = BODY_X + name_w
            by = ry + bi * 36
            add_rect(slide, bx, by + 4, bar_max_w, 26, fill=C["card2"])
            fw = max(4, int(bar_max_w * pct))
            if fw > 4:
                add_rect(slide, bx, by + 4, fw, 26, fill=color)
            add_text(slide, bx + bar_max_w + 8, by + 6, 110, 22,
                     f"{yes}/{n}  ({pct*100:.0f}%)", 11, color, bold=True)


def s_helmet_ok_frames(prs: Presentation, baseline: list[dict], data_dir: Path) -> None:
    rows = baseline
    person = [r for r in rows if r["result"]["detections"]["person"]["count"] > 0]
    ok_frames = [r for r in person
                 if r["result"]["detections"]["helmet"]["count"] > 0
                 and r["result"]["queries"].get("missing_helmet", {}).get("label") == "no"][:6]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1280, 720, fill=C["bg"])
    hdr(slide, "Helmet — YOLO 正確識別的幀（detect=1, missing_helmet=no）",
        f"HELMET OK FRAMES · detect_helmet=1 · query=no · n={len(ok_frames)} SHOWN")

    cols, rows_n, gap = 3, 2, 10
    iw = (BODY_W - gap * (cols - 1)) // cols
    ih = (BODY_H - gap * (rows_n - 1)) // rows_n

    for idx, r in enumerate(ok_frames):
        gx = BODY_X + (idx % cols) * (iw + gap)
        gy = BODY_Y + 8 + (idx // cols) * (ih + gap)
        fb = get_frame_bytes(r["video"], float(r["time_sec"]), data_dir, width=440)
        if fb is None:
            fb = placeholder_img(440, 248)
        add_rect(slide, gx, gy, iw, ih, fill=C["card2"])
        add_img(slide, fb, gx, gy, iw, ih - 46)
        add_rect(slide, gx, gy + ih - 46, iw, 46, fill=C["hdr_bg"])
        add_rect(slide, gx, gy, iw, 4, fill=C["teal"])
        add_text(slide, gx + 5, gy + ih - 43, iw - 10, 22,
                 "detect_helmet=1  missing_helmet=no  ✓ 正確", 10, C["teal"], bold=True)
        add_text(slide, gx + 5, gy + ih - 22, iw - 10, 18,
                 f"{Path(r['video']).name[:30]} @ {r['time_sec']}s", 9, "#cbd5e1")


def s_abnormal_behavior(prs: Presentation, baseline: list[dict], data_dir: Path) -> None:
    rows = baseline
    person = [r for r in rows if r["result"]["detections"]["person"]["count"] > 0]
    ab_yes = [r for r in person
              if r["result"]["queries"].get("abnormal_behavior", {}).get("label") == "yes"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1280, 720, fill=C["bg"])
    hdr(slide, f"abnormal_behavior 偵測結果（n={len(ab_yes)} 幀）",
        "QUERY: abnormal_behavior · 8.8% FP RATE (3/34) · LOW RISK")

    # stat card
    card_rect(slide, BODY_X, BODY_Y + 8, 280, 300, accent_color=C["blue"], accent_side="left")
    add_text(slide, BODY_X + 16, BODY_Y + 20, 250, 70,
             f"{len(ab_yes)}/34", 48, C["blue"], bold=True)
    add_text(slide, BODY_X + 16, BODY_Y + 96, 250, 40,
             "有人幀回答 yes", 13, C["sub"])
    add_text(slide, BODY_X + 16, BODY_Y + 144, 250, 24,
             "FP 率 = 8.8%", 14, C["blue"], bold=True)
    add_text(slide, BODY_X + 16, BODY_Y + 174, 250, 100,
             "3 幀皆為工人彎腰靠近貨車，模型判定為潛在危險姿態", 12, C["sub"], wrap=True)

    # 3 image cards horizontal
    iw, ih = 300, 560
    gap = 14
    total_w = len(ab_yes) * iw + (len(ab_yes) - 1) * gap
    start_x = BODY_X + 298
    for idx, r in enumerate(ab_yes):
        gx = start_x + idx * (iw + gap)
        gy = BODY_Y + 8
        answer = r["result"]["queries"]["abnormal_behavior"].get("answer", "")
        fb = get_frame_bytes(r["video"], float(r["time_sec"]), data_dir, width=400)
        if fb is None:
            fb = placeholder_img(400, 225)
        add_rect(slide, gx, gy, iw, ih, fill=C["card2"])
        add_rect(slide, gx, gy, iw, 4, fill=C["blue"])
        add_img(slide, fb, gx, gy + 4, iw, 300)
        add_rect(slide, gx, gy + 308, iw, 252, fill=C["hdr_bg"])
        add_text(slide, gx + 8, gy + 314, iw - 16, 22,
                 "abnormal_behavior = yes", 11, C["blue"], bold=True)
        add_text(slide, gx + 8, gy + 340, iw - 16, 120,
                 answer[:160], 10, "#cbd5e1", wrap=True)
        add_text(slide, gx + 8, gy + 470, iw - 16, 20,
                 f"{Path(r['video']).name[:32]} @ {r['time_sec']}s", 9, "#64748b")


def s_capability_overview(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1280, 720, fill=C["bg"])
    hdr(slide, "Moondream 整體能力定位",
        "CAPABILITY MATRIX · WHAT WORKS, WHAT DOESN'T, WHAT NEEDS HARDWARE FIX")

    categories = [
        # (label, detect_ok, query_ok, status_color, status_text, note)
        ("missing_helmet",    True,  False, C["amber"], "部分可用",
         "detect 可靠；query 誤報 32.4%。\nhelmet_area_visible 補充後 FP 降到 ~5%。\n→ Pipeline 設計後可部署"),
        ("missing_face_mask", True,  False, C["amber"], "部分可用",
         "detect 可靠；query 誤報 70.6%。\nanything_on_face 降到 26.5%，仍有 9 幀無法確認。\n→ 需補充側面鏡頭"),
        ("abnormal_behavior", False, True,  C["blue"],  "query 可用",
         "無 detect，純 query。FP 8.8%（3/34）。\n偵測到工人彎腰靠近貨車等危險姿態。\n→ 可直接整合，加規則過濾"),
        ("smoking",           False, True,  C["teal"],  "✓ 無誤報",
         "0/34 誤報。本批資料無吸菸案例，\n無法評估 true positive 表現。\n→ 需更多樣本驗證"),
        ("fallen_person",     False, True,  C["teal"],  "✓ 無誤報",
         "0/34 誤報。本批資料無倒地案例。\n→ 需更多樣本驗證"),
    ]

    col_labels = ["Query", "detect()", "query()", "狀態", "說明"]
    col_x = [BODY_X, BODY_X + 230, BODY_X + 340, BODY_X + 460, BODY_X + 600]
    col_w = [220,      100,          110,           130,           648]

    # header row
    add_rect(slide, BODY_X, BODY_Y + 8, BODY_W, 32, fill=C["hdr_bg"])
    for j, (lbl, cx, cw) in enumerate(zip(col_labels, col_x, col_w)):
        add_text(slide, cx + 4, BODY_Y + 14, cw - 8, 20,
                 lbl, 11, "#94a3b8", bold=True)

    row_h = 104
    for i, (label, det_ok, qry_ok, scolor, status, note) in enumerate(categories):
        ry = BODY_Y + 44 + i * row_h
        bg = C["card"] if i % 2 == 0 else C["card2"]
        add_rect(slide, BODY_X, ry, BODY_W, row_h - 2, fill=bg)
        add_rect(slide, BODY_X, ry, 4, row_h - 2, fill=scolor)

        add_text(slide, col_x[0] + 10, ry + 8, col_w[0] - 14, 24,
                 label, 12, C["text"], bold=True)

        for j, ok in enumerate([det_ok, qry_ok]):
            mark, mc = ("✓", C["teal"]) if ok else ("—", C["muted"])
            add_text(slide, col_x[1 + j] + 10, ry + 8, col_w[1 + j] - 14, 24,
                     mark, 16, mc, bold=True, align=PP_ALIGN.CENTER)

        add_rect(slide, col_x[3] + 6, ry + 10, col_w[3] - 12, 26, fill=scolor, radius=True)
        add_text(slide, col_x[3] + 10, ry + 13, col_w[3] - 20, 20,
                 status, 10, "#ffffff", bold=True)

        add_text(slide, col_x[4] + 6, ry + 6, col_w[4] - 12, row_h - 14,
                 note, 11, C["sub"], wrap=True)


def s_summary(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 1280, 720, fill=C["bg"])
    hdr(slide, "實驗結論總覽", "SUMMARY · HELMET + MASK · THREE EXPERIMENT PHASES")

    panels = [
        (C["teal"], "HELMET ✓  Pipeline 已驗證", [
            ("detect_helmet=0 → 問 helmet_area_visible", C["sub"]),
            ("  = no  → 取消告警（遮擋）  6/7", C["teal"]),
            ("  = yes → 升級告警", C["sub"]),
            ("detect_helmet>0 → 信任 detect，query 只做補充", C["sub"]),
        ], "32.4%", "原始 FP 率", "~5%", "Pipeline 後估計"),
        (C["amber"], "MASK ⚠  部分改善，仍有限制", [
            ("face_area_visible → 無效（兩組無差異）", C["sub"]),
            ("anything_on_face → 有效", C["sub"]),
            ("  = yes → 告警取消  15/24", C["teal"]),
            ("  = no  → 告警保留   9/24", C["sub"]),
        ], "70.6%", "原始告警率", "26.5%", "anything_on_face 後"),
    ]

    for i, (accent, title, lines, n1, l1, n2, l2) in enumerate(panels):
        cx = BODY_X + i * 636
        cw = 620
        card_rect(slide, cx, BODY_Y + 8, cw, 560, accent_color=accent)
        add_text(slide, cx + 14, BODY_Y + 20, cw - 28, 30,
                 title, 16, accent, bold=True)
        for li, (txt, col) in enumerate(lines):
            add_text(slide, cx + 14, BODY_Y + 60 + li * 36, cw - 28, 32,
                     txt, 13, col, wrap=False)
        # stats row
        add_rect(slide, cx + 14, BODY_Y + 430, cw - 28, 2, fill=C["border"])
        for si, (num, lbl) in enumerate([(n1, l1), (n2, l2)]):
            sx = cx + 14 + si * ((cw - 28) // 2)
            add_text(slide, sx, BODY_Y + 446, (cw - 28) // 2 - 8, 38,
                     num, 24, accent, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, sx, BODY_Y + 488, (cw - 28) // 2 - 8, 28,
                     lbl, 12, C["muted"], align=PP_ALIGN.CENTER)


# ── Main ───────────────────────────────────────────────────────────────────
def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--run", required=True)
    parser.add_argument("--data", default=str(PROJECT_ROOT / "data"))
    args = parser.parse_args()

    data_dir = Path(args.data)
    run      = args.run
    out_path = OUTPUTS_DIR / run / "presentation_editable.pptx"

    prs = Presentation()
    prs.slide_width  = Emu(9144000)   # 10 inches
    prs.slide_height = Emu(5143500)   # 5.625 inches (16:9)

    print("[data] loading...")
    baseline  = _load_baseline(run)
    phase2    = _load_phase2(run)
    mask_data = _load_mask_data(run)
    print(f"[data] baseline={len(baseline)} phase2={len(phase2)} mask_phase1={len(mask_data[0])}")

    slides = [
        (s_cover,                   (prs, run),                          "cover"),
        (s_fp_overview,             (prs, run),                          "fp overview"),
        (s_per_video_breakdown,     (prs, baseline),                     "per-video breakdown"),
        (s_helmet_fp_causes,        (prs, phase2, data_dir),             "helmet fp causes"),
        (s_helmet_contradiction,    (prs, phase2, data_dir),             "helmet contradiction"),
        (s_helmet_ok_frames,        (prs, baseline, data_dir),           "helmet ok frames"),
        (s_mask_problem,            (prs, mask_data),                    "mask problem"),
        (s_mask_fp_photos,          (prs, mask_data, data_dir),          "mask fp photos"),
        (s_mask_improvement_stats,  (prs, mask_data),                    "mask stats"),
        (s_mask_detected_frames,    (prs, mask_data, data_dir),          "mask detected"),
        (s_mask_missed_frames,      (prs, mask_data, data_dir),          "mask missed"),
        (s_abnormal_behavior,       (prs, baseline, data_dir),           "abnormal behavior"),
        (s_capability_overview,     (prs,),                              "capability overview"),
        (s_summary,                 (prs,),                              "summary"),
    ]
    n = len(slides)
    print(f"[pptx] building {n} slides...")
    for i, (fn, args, label) in enumerate(slides, 1):
        fn(*args)
        print(f"[pptx] slide {i}/{n} {label}")

    prs.save(str(out_path))
    print(f"\n完成：{out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
